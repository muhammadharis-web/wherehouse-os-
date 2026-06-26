from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, items, left, top, width, height, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
    return txBox

def add_accent_line(slide, left, top, width=0.8, height=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_card(slide, text, left, top, width, height, bg_color=RGBColor(0x1E, 0x29, 0x3B)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 1 - Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_shape_bg(slide, RGBColor(0x0A, 0x0F, 0x1E))
add_shape_bg(slide, ACCENT, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

add_text_box(slide, "WAREHOUSE OS", 1.5, 1.5, 10, 1.2, 54, WHITE, True, PP_ALIGN.CENTER)
add_accent_line(slide, 5.8, 2.9, 1.6)
add_text_box(slide, "Multi-Agent Order Fulfillment System", 1.5, 3.3, 10, 0.8, 28, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, "AI-Powered Warehouse Management with Autonomous Agents", 1.5, 4.0, 10, 0.6, 18, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, "Class Presentation - Final Project", 1.5, 5.5, 10, 0.5, 18, RGBColor(0x6B, 0x72, 0x8B), False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 - Project Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, LIGHT_BG)
add_shape_bg(slide, DARK, Inches(0), Inches(0), Inches(5.5), Inches(7.5))

add_text_box(slide, "PROJECT OVERVIEW", 0.8, 0.8, 4.5, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.5, 1.2)
add_text_box(slide, "An intelligent warehouse management system powered by AI agents that automates the entire order fulfillment lifecycle — from order placement to delivery tracking.", 0.8, 2.0, 4.2, 2.5, 14, GRAY)

items = [
    "Real-time order processing & tracking",
    "Multi-agent AI orchestration",
    "Vector database for semantic search",
    "Automated rerouting & delay prediction",
    "SMS/Email notifications (Twilio + SendGrid)",
    "3D warehouse visualization (Three.js)",
    "Smooth animations (Framer Motion + GSAP)",
]
add_bullet_slide(slide, items, 0.8, 4.2, 4.5, 4.5, 14, RGBColor(0xCB, 0xD5, 0xE1))

add_text_box(slide, "TECH STACK OVERVIEW", 6.2, 0.8, 6.5, 0.6, 24, DARK, True)
add_accent_line(slide, 6.2, 1.5, 1.2)

cards = [
    ("Backend", "Python 3.12 • FastAPI\nSQLAlchemy • Celery\nRedis • Qdrant • OpenAI", 6.2, 2.0, 3.0, 1.6),
    ("Frontend", "Next.js 16 • React 19\nTypeScript • Tailwind CSS v4\nFramer Motion • Three.js", 9.5, 2.0, 3.0, 1.6),
    ("Database", "PostgreSQL 16 (Prod)\nSQLite (Dev)\nQdrant Vector DB\nRedis Cache", 6.2, 3.9, 3.0, 1.6),
    ("AI / Agents", "OpenAI GPT-4o\nOpenAI Agents SDK\n5 Autonomous Agents\nVector Embeddings", 9.5, 3.9, 3.0, 1.6),
    ("DevOps", "Docker • docker-compose\nGit • Ruff • mypy\npytest • Vitest\nESLint", 6.2, 5.8, 3.0, 1.6),
    ("Notifications", "Twilio (SMS)\nSendGrid (Email)\nEasyPost (Shipping)\nSmartyStreets (Addr)", 9.5, 5.8, 3.0, 1.6),
]
for text, desc, left, top, w, h in cards:
    add_card(slide, f"{text}\n—————————\n{desc}", left, top, w, h)

# ============================================================
# SLIDE 3 - System Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)

add_text_box(slide, "SYSTEM ARCHITECTURE", 0.8, 0.4, 8, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.1, 1.2)

add_card(slide, "Frontend\nNext.js 16 • React 19\nTailwind CSS v4\nPort 3000", 0.5, 1.6, 2.8, 1.5, RGBColor(0x1E, 0x40, 0x73))
add_card(slide, "Backend API\nFastAPI • Uvicorn\nREST + WebSocket\nPort 8000", 3.6, 1.6, 2.8, 1.5, RGBColor(0x1E, 0x40, 0x73))
add_card(slide, "AI Agents\nOpenAI Agents SDK\n5 Specialized Agents\nGPT-4o Powered", 6.7, 1.6, 2.8, 1.5, RGBColor(0x7C, 0x3A, 0xED))
add_card(slide, "Vector DB\nQdrant\n4 Collections\n1536-dim Vectors", 9.8, 1.6, 2.8, 1.5, RGBColor(0x7C, 0x3A, 0xED))

add_card(slide, "PostgreSQL 16\nRelational DB\nOrders, Shipments\nFulfillment Centers", 0.5, 3.5, 2.8, 1.5, RGBColor(0x06, 0x5F, 0x46))
add_card(slide, "Redis 7\nMessage Broker\nCelery Backend\nCache Layer", 3.6, 3.5, 2.8, 1.5, RGBColor(0x06, 0x5F, 0x46))
add_card(slide, "Celery Workers\nAsync Tasks\nMonitor Cycles\nBackground Jobs", 6.7, 3.5, 2.8, 1.5, RGBColor(0x06, 0x5F, 0x46))
add_card(slide, "External APIs\nTwilio • SendGrid\nOpenAI • EasyPost\nSmartyStreets", 9.8, 3.5, 2.8, 1.5, RGBColor(0x78, 0x1A, 0x0F))

add_text_box(slide, "Data Flow:  User → Frontend (Next.js) → API (FastAPI) → AI Agents → DB/External Services → Response", 0.5, 5.6, 12.5, 0.8, 14, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, "Deployment: Docker Containers (PostgreSQL + Redis + API + Celery Workers + Celery Beat)", 0.5, 6.3, 12.5, 0.8, 14, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4 - AI Agents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, LIGHT_BG)
add_shape_bg(slide, DARK, Inches(0), Inches(0), Inches(13.333), Inches(1.3))

add_text_box(slide, "AI AGENTS ARCHITECTURE", 0.8, 0.35, 10, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.05, 1.2)

agents = [
    ("Fulfillment Orchestrator", "Central coordinator that runs the full monitor cycle: checks delays, evaluates reroutes, predicts failures, and optimizes costs across all agents.", ACCENT),
    ("Routing Agent", "Selects optimal fulfillment center & carrier rate for new orders based on cost, location, and SLA requirements.", GREEN),
    ("Monitor Agent", "Queries active shipments every N seconds, detects delays and overdue status, and triggers alerts.", ORANGE),
    ("Rerouting Agent", "Evaluates alternative carriers when delays occur, checks feasibility, and executes automated reroutes.", RED),
    ("Communication Agent", "Sends delay alerts via email (SendGrid) and SMS (Twilio) with customer-preferred channels.", ACCENT),
    ("Prediction Agent", "Predicts failure probability for shipments based on delay history, carrier statistics, and risk factors using vector similarity.", GREEN),
    ("Cost Optimizer", "Analyzes shipping costs per monitor cycle, identifies trends, and generates cost-reduction recommendations.", ORANGE),
]

for i, (name, desc, color) in enumerate(agents):
    y = 1.7 + i * 0.78
    add_shape_bg(slide, color, Inches(0.5), Inches(y), Inches(0.08), Inches(0.6))
    add_text_box(slide, name, 0.8, y - 0.02, 3.5, 0.4, 16, DARK, True)
    add_text_box(slide, desc, 4.5, y - 0.02, 8.2, 0.6, 12, RGBColor(0x47, 0x4E, 0x63))

# ============================================================
# SLIDE 5 - Guardrails
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)

add_text_box(slide, "BUSINESS GUARDRAILS", 0.8, 0.5, 8, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.2, 1.2)

add_text_box(slide, "Safety rules enforced during every monitor cycle to ensure business compliance:", 0.8, 1.5, 11, 0.5, 14, GRAY)

guardrails = [
    ("SLA Compliance", "Flags shipments exceeding critical SLA hours thresholds"),
    ("Cost Cap", "Prevents reroutes exceeding maximum allowed cost increase percentage"),
    ("Notification Frequency", "Limits notifications per order to avoid spam"),
    ("Failed Delivery Threshold", "Flags shipments exceeding overdue delivery threshold"),
    ("Carrier Diversity", "Prevents switching to same/similar monopoly carriers"),
    ("Address Validation", "Validates address format, ZIP code, and street indicators"),
]

for i, (name, desc) in enumerate(guardrails):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.3 + row * 2.3
    add_card(slide, f"🛡  {name}\n{desc}", x, y, 3.6, 1.8, RGBColor(0x1E, 0x29, 0x3B))

add_text_box(slide, "⚡ All guardrails run asynchronously during each Celery Beat monitor cycle (every 15 minutes)", 0.8, 6.6, 11, 0.5, 13, GRAY)

# ============================================================
# SLIDE 6 - Frontend Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, LIGHT_BG)
add_shape_bg(slide, DARK, Inches(0), Inches(0), Inches(13.333), Inches(1.3))

add_text_box(slide, "FRONTEND HIGHLIGHTS", 0.8, 0.35, 10, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.05, 1.2)

features = [
    ("Landing Page", "3D Warehouse Scene (Three.js)\nHero animations (Framer Motion, GSAP)\nSmooth scrolling (Lenis)\nParallax layers, Tilt cards", 0.5, 1.6),
    ("Dashboard", "Real-time order tracking\nInventory zones panel\nAnalytics & metrics\nAgent activity logs", 3.5, 1.6),
    ("AI Assistant", "Chat interface (OpenAI GPT-4o-mini)\nStreaming responses (Vercel AI SDK)\nSmart suggestions & insights\nContext-aware replies", 6.5, 1.6),
    ("UI/UX", "Radix UI accessible components\nTailwind CSS v4 styling\nFramer Motion animations\nGSAP advanced timelines\nLucide icons • Dark theme", 9.5, 1.6),
    ("Auth & Forms", "NextAuth.js authentication\nGoogle + Facebook OAuth\nReact Hook Form + Zod v4\nJWT session management", 0.5, 4.2),
    ("Visual Effects", "Noise overlay • Cursor follower\nMagnetic buttons • Rotating borders\nText reveal animations\nAnimated counters • Scroll animations", 3.5, 4.2),
    ("Pages", "/dashboard • /orders • /agents\n/analytics • /settings\n/workflows • /monitoring\n/team", 6.5, 4.2),
    ("Testing", "Vitest • React Testing Library\nJest DOM matchers\nUser event simulation\njsdom environment", 9.5, 4.2),
]
for name, desc, left, top in features:
    add_card(slide, f"{name}\n—————————\n{desc}", left, top, 2.8, 2.3, WHITE)
    # Fix text color for white cards
    for shape in slide.shapes:
        if shape == slide.shapes[-1]:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.rgb == WHITE:
                            run.font.color.rgb = DARK
                    except:
                        run.font.color.rgb = DARK

# ============================================================
# SLIDE 7 - Tech Stack Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)

add_text_box(slide, "TECHNOLOGY STACK", 0.8, 0.4, 8, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.1, 1.2)

stacks = [
    ("Languages & Runtimes", "Python 3.12, TypeScript 5, Node.js 20+, pnpm 9+"),
    ("Backend Framework", "FastAPI, Uvicorn, Pydantic, SQLAlchemy 2.0 (async)"),
    ("Frontend Framework", "Next.js 16 (App Router), React 19, Tailwind CSS v4"),
    ("Database & Cache", "PostgreSQL 16, SQLite (dev), Redis 7, Qdrant Vector DB"),
    ("AI / LLM", "OpenAI GPT-4o, OpenAI Agents SDK, text-embedding-3-small"),
    ("Task Queue", "Celery, Celery Beat, Redis (broker + backend)"),
    ("Auth", "NextAuth.js, JWT (python-jose), OAuth2, bcrypt"),
    ("UI Components", "Radix UI, lucide-react, class-variance-authority, tailwind-merge"),
    ("Animations", "Framer Motion 12, GSAP 3, Three.js, Lenis smooth scroll"),
    ("Forms & Validation", "React Hook Form, Zod v4, @hookform/resolvers"),
    ("Notifications", "Twilio (SMS), SendGrid (Email)"),
    ("External APIs", "OpenAI, EasyPost (Shipping), SmartyStreets (Address)"),
    ("DevOps & Tools", "Docker, docker-compose, Git, ESLint, Ruff, mypy, pytest, Vitest"),
    ("Testing", "pytest-asyncio, Vitest, React Testing Library, Locust (load test)"),
]

for i, (name, desc) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.4 + row * 0.42
    add_text_box(slide, f"▸ {name}", x, y, 2.8, 0.35, 12, ACCENT, True)
    add_text_box(slide, desc, x + 2.8, y, 3.4, 0.35, 11, GRAY)

# ============================================================
# SLIDE 8 - Database Schema
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, LIGHT_BG)
add_shape_bg(slide, DARK, Inches(0), Inches(0), Inches(13.333), Inches(1.3))

add_text_box(slide, "DATABASE & VECTOR STORE", 0.8, 0.35, 10, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.05, 1.2)

add_text_box(slide, "Relational Database (PostgreSQL / SQLite)", 0.5, 1.6, 6, 0.4, 18, DARK, True)
add_card(slide, "Orders\n• Order ID, Customer, Items\n• Status, Total, Timeline\n• Agent Assignments", 0.5, 2.2, 2.8, 1.8, RGBColor(0x1E, 0x40, 0x73))
add_card(slide, "Shipments\n• Shipment ID, Carrier\n• Tracking, Status\n• Delays, Reroutes", 3.6, 2.2, 2.8, 1.8, RGBColor(0x1E, 0x40, 0x73))
add_card(slide, "Fulfillment Centers\n• Location, Capacity\n• Active Hours\n• Carrier Partners", 6.7, 2.2, 2.8, 1.8, RGBColor(0x1E, 0x40, 0x73))
add_card(slide, "Carrier Rates\n• Rate Tables\n• Service Levels\n• Cost per Zone", 9.8, 2.2, 2.8, 1.8, RGBColor(0x1E, 0x40, 0x73))

add_text_box(slide, "Vector Database (Qdrant) - 1536-dimension embeddings", 0.5, 4.3, 6, 0.4, 18, DARK, True)
add_card(slide, "shipment_events\nVector embeddings of\nshipment tracking events\nfor similarity search", 0.5, 4.9, 2.8, 1.6, RGBColor(0x7C, 0x3A, 0xED))
add_card(slide, "product_catalog\nProduct descriptions &\nfeatures as vectors\nfor semantic search", 3.6, 4.9, 2.8, 1.6, RGBColor(0x7C, 0x3A, 0xED))
add_card(slide, "customer_order_history\nCustomer patterns &\npreferences for\npersonalized routing", 6.7, 4.9, 2.8, 1.6, RGBColor(0x7C, 0x3A, 0xED))
add_card(slide, "agent_decisions\nAgent reasoning traces\nfor audit & continuous\nimprovement", 9.8, 4.9, 2.8, 1.6, RGBColor(0x7C, 0x3A, 0xED))

add_text_box(slide, "⚡ Mock fallback available — Qdrant not required for local development", 0.5, 6.8, 12, 0.4, 12, RGBColor(0x6B, 0x72, 0x8B))

# ============================================================
# SLIDE 9 - Agent Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)

add_text_box(slide, "ORDER FULFILLMENT WORKFLOW", 0.8, 0.4, 10, 0.6, 28, WHITE, True)
add_accent_line(slide, 0.8, 1.1, 1.2)

steps = [
    ("1", "Order Placed", "Customer places order\nvia frontend UI", ACCENT),
    ("2", "Routing Agent", "Selects best fulfillment\ncenter & carrier rate", GREEN),
    ("3", "Order Processing", "Order assigned to FC\nInventory reserved", ORANGE),
    ("4", "Fulfillment", "Items picked & packed\nShipping label generated", RED),
    ("5", "In Transit", "Real-time tracking\nMonitor agent active", ACCENT),
    ("6", "Delay Detection", "Delay predicted via\nvector similarity search", GREEN),
    ("7", "Rerouting", "Alternative carrier\nselected & executed", ORANGE),
    ("8", "Communication", "SMS/Email alerts sent\nvia Twilio & SendGrid", RED),
    ("9", "Delivery", "Order delivered\nCustomer notified", GREEN),
]

for i, (num, name, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.5 + row * 1.9
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    
    add_text_box(slide, name, x, y + 0.7, 3.8, 0.4, 15, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, desc, x, y + 1.1, 3.8, 0.6, 11, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10 - Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)
add_shape_bg(slide, RGBColor(0x0A, 0x0F, 0x1E))
add_shape_bg(slide, ACCENT, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

add_text_box(slide, "THANK YOU", 1.5, 2.0, 10, 1.2, 48, WHITE, True, PP_ALIGN.CENTER)
add_accent_line(slide, 5.5, 3.3, 2.2)
add_text_box(slide, "Warehouse OS — Multi-Agent Order Fulfillment System", 1.5, 3.6, 10, 0.6, 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, "GitHub: github.com/muhammadharis-web/wherehouse-os-", 1.5, 4.5, 10, 0.5, 16, ACCENT, False, PP_ALIGN.CENTER)
add_text_box(slide, "Questions?", 1.5, 5.5, 10, 1, 32, GRAY, False, PP_ALIGN.CENTER)

# Save
output_path = "C:\\Users\\AC\\Desktop\\final project\\Warehouse_OS_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
